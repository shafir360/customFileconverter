import os
from flask import Flask, jsonify, request
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

app = Flask(__name__)

# Set up upload configuration
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pptx'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def allowed_file(filename):
    """Check if the uploaded file has an allowed extension (pptx)."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pptx_improved(file_path):
    """
    Extracts text from a PPTX file with improved formatting.
    
    This function goes beyond a simple text grab:
      - It processes each slide and adds a header (e.g. "Slide 1:")
      - For each shape with a text frame, it iterates over all paragraphs and runs
        to reconstruct the text while preserving basic formatting.
      - It handles table shapes by iterating over table rows and cells.
    
    :param file_path: Path to the PPTX file.
    :return: A string that contains the structured text extracted from the presentation.
    """
    prs = Presentation(file_path)
    slides_text = []
    
    # Process each slide, numbering them for clarity.
    for idx, slide in enumerate(prs.slides, start=1):
        slide_content = [f"Slide {idx}:"]
        for shape in slide.shapes:
            # If the shape has a text_frame attribute, process it.
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                shape_text_lines = []
                for paragraph in shape.text_frame.paragraphs:
                    # Concatenate all runs in the paragraph to form a complete line.
                    paragraph_text = "".join(run.text for run in paragraph.runs)
                    if paragraph_text.strip():
                        shape_text_lines.append(paragraph_text.strip())
                if shape_text_lines:
                    slide_content.append("\n".join(shape_text_lines))
            # If the shape is a table, extract text from its cells.
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                table_text = []
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_cells:
                        table_text.append(" | ".join(row_cells))
                if table_text:
                    slide_content.append("\n".join(table_text))
        slides_text.append("\n".join(slide_content))
    return "\n\n".join(slides_text)

@app.route('/extract-text', methods=['POST'])
def extract_text():
    """
    Endpoint to extract text from an uploaded PPTX file.

    Expects a form-data POST request with the file attached under the key 'file'.
    """
    if 'file' not in request.files:
        return jsonify({'error': 'No file part in the request'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)

        # Ensure the upload directory exists.
        os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
        file.save(filepath)

        try:
            # Use the improved extraction function.
            extracted_text = extract_text_from_pptx_improved(filepath)
        except Exception as e:
            os.remove(filepath)
            return jsonify({'error': f'Error processing PPTX file: {str(e)}'}), 500

        # Clean up the temporary file.
        os.remove(filepath)
        return jsonify({'extracted_text': extracted_text})
    else:
        return jsonify({'error': 'Invalid file type. Only PPTX files are allowed.'}), 400

@app.route('/')
def index():
    return jsonify({"Choo Choo": "Welcome to your Flask app 🚅"})

if __name__ == '__main__':
    app.run(debug=True, port=os.getenv("PORT", default=5000))
