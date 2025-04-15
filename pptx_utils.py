import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ALLOWED_EXTENSIONS = {'pptx'}

def allowed_file(filename):
    """Check if the uploaded file has an allowed extension (pptx)."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pptx_markdown(file_path):
    """
    Extracts text from a PPTX file and returns it in Markdown format.
    """
    prs = Presentation(file_path)
    slides_md = []
    slide_num = 1  # Continuous numbering for non-empty slides

    for slide in prs.slides:
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = "".join(run.text for run in paragraph.runs).strip()
                    if paragraph_text:
                        slide_lines.append(paragraph_text)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_lines.append(row_text)
        if slide_lines:
            header = f"## Slide {slide_num}"
            slide_md = "\n".join([header, ""] + slide_lines)
            slides_md.append(slide_md)
            slide_num += 1

    return "\n\n".join(slides_md)
