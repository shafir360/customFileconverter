from flask import Flask, request, jsonify
from pptx import Presentation
import tempfile
import os

app = Flask(__name__)

# Existing route (from the template)
@app.route("/")
def index():
    return "Hello, World!"

# Function to extract text from a PowerPoint file
def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

# New /extract endpoint
@app.route("/extract", methods=["POST"])
def extract():
    # Ensure a file was uploaded under the "file" key
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    try:
        # Save file to a temporary location
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            file.save(tmp.name)
            # Extract text from the temporary file
            extracted_text = extract_text_from_pptx(tmp.name)
            # Remove the temporary file after processing
            os.unlink(tmp.name)
        return jsonify({"text": extracted_text})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    # Railway usually expects the app to be available on 0.0.0.0
    app.run(host="0.0.0.0", port=8000)
