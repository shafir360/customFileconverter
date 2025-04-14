from flask import Flask, request, jsonify
from pptx import Presentation
import tempfile
import os

app = Flask(__name__)

def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

@app.route('/extract', methods=['POST'])
def extract():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    # Save the file to a temporary location
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        file.save(tmp.name)
        extracted_text = extract_text_from_pptx(tmp.name)
        # Clean up the temporary file after extraction
        os.unlink(tmp.name)
    
    return jsonify({"text": extracted_text})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8000)
